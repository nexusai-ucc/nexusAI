"""Extracción de texto de documentos: PDF, DOCX, TXT, PPTX, XLSX, CSV, MD, HTML — CONT-03.

Limitaciones:
  - PDF: texto seleccionable; si una página no tiene texto se intenta OCR (Tesseract)
    sobre su render como imagen. Lanza ValueError solo si ni el texto ni el OCR
    producen contenido.
  - DOCX: extrae párrafos y celdas de tablas. No extrae headers/footers ni imágenes.
  - PPTX: extrae texto de shapes y tablas; si una slide no tiene texto pero sí
    imágenes, se intenta OCR sobre esas imágenes.
  - XLSX: extrae el valor de todas las celdas no vacías de todas las hojas.
  - CSV / TXT: UTF-8 con fallback a latin-1.
  - MD: se trata como texto plano (sin parseo de sintaxis Markdown).
  - HTML: se extrae el texto visible, descartando tags y scripts/estilos.
"""

from __future__ import annotations

import logging
from io import BytesIO

import pdfplumber


logger = logging.getLogger("nexusai.documents.extractor")

_MIME_PDF = "application/pdf"
_MIME_DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_MIME_TXT = "text/plain"
_MIME_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_MIME_XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_MIME_CSV = "text/csv"
_MIME_MD = "text/markdown"
_MIME_HTML = "text/html"

SUPPORTED_MIME_TYPES = {
    _MIME_PDF,
    _MIME_DOCX,
    _MIME_TXT,
    _MIME_PPTX,
    _MIME_XLSX,
    _MIME_CSV,
    _MIME_MD,
    _MIME_HTML,
}


def extract_text(file_bytes: bytes, mime_type: str = _MIME_PDF) -> str:
    """Extrae texto de un documento según su mime_type.

    Args:
        file_bytes: Contenido binario del archivo.
        mime_type: MIME type del archivo.

    Returns:
        Texto extraído como string, listo para el pipeline de chunking.

    Raises:
        ValueError: si el archivo está vacío, corrupto, no tiene texto extraíble,
                    o el mime_type no está soportado.
    """
    if not file_bytes:
        raise ValueError("El archivo está vacío o no contiene datos válidos.")

    if mime_type == _MIME_PDF:
        return _extract_pdf(file_bytes)
    elif mime_type == _MIME_DOCX:
        return _extract_docx(file_bytes)
    elif mime_type == _MIME_TXT:
        return _extract_txt(file_bytes)
    elif mime_type == _MIME_PPTX:
        return _extract_pptx(file_bytes)
    elif mime_type == _MIME_XLSX:
        return _extract_xlsx(file_bytes)
    elif mime_type == _MIME_CSV:
        return _extract_csv(file_bytes)
    elif mime_type == _MIME_MD:
        return _extract_txt(file_bytes)
    elif mime_type == _MIME_HTML:
        return _extract_html(file_bytes)
    else:
        raise ValueError(
            f"Tipo de archivo no soportado: {mime_type!r}. "
            f"Tipos aceptados: {sorted(SUPPORTED_MIME_TYPES)}"
        )


def _ocr_image(image_bytes: bytes) -> str:
    """Aplica OCR (Tesseract) sobre una imagen y devuelve el texto detectado.

    Importación lazy de pytesseract/Pillow: son dependencias pesadas (requieren el
    binario de sistema tesseract-ocr) que solo se necesitan en el camino OCR.
    """
    try:
        import pytesseract
        from PIL import Image

        img = Image.open(BytesIO(image_bytes))
        return pytesseract.image_to_string(img, lang="spa+eng").strip()
    except Exception as exc:
        logger.warning("OCR falló sobre una imagen: %s", exc)
        return ""


def _extract_pdf(file_bytes: bytes) -> str:
    try:
        with pdfplumber.open(BytesIO(file_bytes)) as pdf:
            pages_text = []
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    pages_text.append(page_text.strip())
    except Exception as exc:
        raise ValueError(
            "No se pudo leer el PDF. El archivo puede estar corrupto o no ser un PDF válido."
        ) from exc

    if not pages_text:
        logger.info("PDF sin texto seleccionable, intentando OCR.")
        pages_text = _ocr_pdf(file_bytes)

    if not pages_text:
        raise ValueError(
            "El PDF no tiene texto extraíble, ni siquiera vía OCR. "
            "Verificá que el archivo no esté dañado o en blanco."
        )

    return "\n\n".join(pages_text).strip()


def _ocr_pdf(file_bytes: bytes) -> list[str]:
    """Renderiza cada página de un PDF como imagen y aplica OCR. Importación lazy."""
    try:
        from pdf2image import convert_from_bytes

        images = convert_from_bytes(file_bytes)
    except Exception as exc:
        logger.warning("No se pudo renderizar el PDF para OCR: %s", exc)
        return []

    pages_text: list[str] = []
    for image in images:
        buf = BytesIO()
        image.save(buf, format="PNG")
        text = _ocr_image(buf.getvalue())
        if text:
            pages_text.append(text)
    return pages_text


def _extract_docx(file_bytes: bytes) -> str:
    try:
        from docx import Document  # python-docx — importación lazy para no afectar startup
        doc = Document(BytesIO(file_bytes))
    except Exception as exc:
        raise ValueError(
            "No se pudo leer el DOCX. El archivo puede estar corrupto o no ser un DOCX válido."
        ) from exc

    parts: list[str] = []

    # Párrafos del cuerpo principal
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    # Celdas de todas las tablas
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    parts.append(text)

    if not parts:
        raise ValueError(
            "El DOCX no tiene texto extraíble. Verificá que el archivo tenga contenido de texto."
        )

    return "\n\n".join(parts).strip()


def _extract_pptx(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation  # python-pptx — importación lazy
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(BytesIO(file_bytes))
    except Exception as exc:
        raise ValueError(
            "No se pudo leer el PPTX. El archivo puede estar corrupto o no ser un PPTX válido."
        ) from exc

    parts: list[str] = []

    for slide in prs.slides:
        slide_parts: list[str] = []
        image_shapes = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_parts.append(text)
            elif getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            slide_parts.append(text)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_shapes.append(shape)

        # Solo recurrimos a OCR si la slide no aportó texto por otra vía:
        # evita procesar logos/decoraciones en slides que ya tienen contenido.
        if not slide_parts:
            for shape in image_shapes:
                ocr_text = _ocr_image(shape.image.blob)
                if ocr_text:
                    slide_parts.append(ocr_text)

        parts.extend(slide_parts)

    if not parts:
        raise ValueError(
            "El PPTX no tiene texto extraíble, ni siquiera vía OCR de las imágenes."
        )

    return "\n\n".join(parts).strip()


def _extract_xlsx(file_bytes: bytes) -> str:
    try:
        import openpyxl  # importación lazy

        workbook = openpyxl.load_workbook(
            BytesIO(file_bytes), read_only=True, data_only=True
        )
    except Exception as exc:
        raise ValueError(
            "No se pudo leer el XLSX. El archivo puede estar corrupto o no ser un XLSX válido."
        ) from exc

    parts: list[str] = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(
                str(cell).strip()
                for cell in row
                if cell is not None and str(cell).strip()
            )
            if row_text:
                parts.append(row_text)

    if not parts:
        raise ValueError(
            "El XLSX no tiene datos extraíbles. Verificá que las hojas tengan contenido."
        )

    return "\n".join(parts).strip()


def _extract_csv(file_bytes: bytes) -> str:
    import csv

    text = _decode_text(file_bytes)

    parts: list[str] = []
    for row in csv.reader(text.splitlines()):
        row_text = " | ".join(cell.strip() for cell in row if cell.strip())
        if row_text:
            parts.append(row_text)

    if not parts:
        raise ValueError(
            "El CSV no tiene datos extraíbles. Verificá que el archivo tenga contenido."
        )

    return "\n".join(parts).strip()


def _extract_html(file_bytes: bytes) -> str:
    try:
        from bs4 import BeautifulSoup  # importación lazy

        text = _decode_text(file_bytes)
        soup = BeautifulSoup(text, "html.parser")

        for tag in soup(["script", "style"]):
            tag.decompose()

        raw_text = soup.get_text(separator="\n")
    except Exception as exc:
        raise ValueError(
            "No se pudo leer el HTML. El archivo puede estar corrupto o mal formado."
        ) from exc

    lines = [line.strip() for line in raw_text.splitlines() if line.strip()]
    if not lines:
        raise ValueError(
            "El HTML no tiene texto extraíble. Verificá que el archivo tenga contenido visible."
        )

    return "\n".join(lines).strip()


def _extract_txt(file_bytes: bytes) -> str:
    text = _decode_text(file_bytes)
    if not text:
        raise ValueError(
            "El archivo de texto no pudo decodificarse en UTF-8 ni latin-1."
        )
    return text


def _decode_text(file_bytes: bytes) -> str:
    """Decodifica bytes a texto probando UTF-8 y luego latin-1 como fallback."""
    for encoding in ("utf-8", "latin-1"):
        try:
            text = file_bytes.decode(encoding).strip()
            if text:
                return text
        except UnicodeDecodeError:
            continue
    return ""
