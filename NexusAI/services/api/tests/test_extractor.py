"""
Tests del extractor de documentos — la pieza que convierte bytes de un
archivo (PDF, DOCX, PPTX, XLSX, CSV, MD, HTML, TXT) en texto plano.

Generamos los archivos in-memory con las mismas librerías que el producto usa
para leerlos (reportlab, python-pptx, openpyxl) para no depender de fixtures
en disco. Si alguna no está instalada, los tests correspondientes se saltean
con un mensaje claro (no son dependencias del producto, solo de tests).
"""

from __future__ import annotations

from io import BytesIO
from unittest import mock

import pytest

from app.documents.extractor import extract_text


PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
CSV_MIME = "text/csv"
MD_MIME = "text/markdown"
HTML_MIME = "text/html"


# ============================================================
# Helpers — generar PDFs in-memory
# ============================================================

def _make_pdf_bytes(text_per_page: list[str]) -> bytes:
    """
    Crea un PDF con N páginas, una página por entrada de `text_per_page`.
    Devuelve los bytes serializados.
    """
    reportlab = pytest.importorskip(
        "reportlab",
        reason="reportlab no instalado — pip install reportlab para correr estos tests",
    )
    from reportlab.pdfgen import canvas  # type: ignore[import-not-found]
    from reportlab.lib.pagesizes import A4  # type: ignore[import-not-found]

    buf = BytesIO()
    c = canvas.Canvas(buf, pagesize=A4)
    for page_text in text_per_page:
        # Escribir el texto a una altura razonable. Si hay líneas largas
        # las cortamos manualmente para que pdfplumber las lea bien.
        y = 800
        for line in page_text.split("\n"):
            c.drawString(50, y, line)
            y -= 20
        c.showPage()
    c.save()
    return buf.getvalue()


# ============================================================
# Happy path
# ============================================================

def test_extract_single_page_pdf():
    pdf_bytes = _make_pdf_bytes(["Una derivada mide la tasa instantánea."])
    text = extract_text(pdf_bytes)
    assert "derivada" in text.lower()
    assert "tasa" in text.lower()


def test_extract_multi_page_pdf_joins_pages():
    pdf_bytes = _make_pdf_bytes([
        "Pagina uno con contenido inicial.",
        "Pagina dos con mas contenido.",
        "Pagina tres ultima.",
    ])
    text = extract_text(pdf_bytes)
    assert "uno" in text.lower()
    assert "dos" in text.lower()
    assert "tres" in text.lower()
    # Las páginas se separan con doble salto de línea.
    assert "\n\n" in text


def test_extract_returns_stripped_text():
    """No queremos whitespace inicial/final en la salida."""
    pdf_bytes = _make_pdf_bytes(["contenido"])
    text = extract_text(pdf_bytes)
    assert text == text.strip()


# ============================================================
# Validaciones de input
# ============================================================

def test_empty_bytes_raises():
    with pytest.raises(ValueError, match="vacío"):
        extract_text(b"")


def test_garbage_bytes_raises():
    """Bytes que no son un PDF válido → ValueError descriptivo."""
    with pytest.raises(ValueError, match=r"PDF|corrupto"):
        extract_text(b"esto no es un PDF, son bytes random")


def test_pdf_with_no_extractable_text_raises():
    """
    Un PDF sin texto (ej. solo imagen escaneada) y sin OCR disponible/exitoso
    → ValueError. Para simularlo creamos un PDF con una página vacía (sin
    drawString) y forzamos que el fallback OCR no encuentre texto.
    """
    pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas  # type: ignore[import-not-found]

    buf = BytesIO()
    c = canvas.Canvas(buf)
    c.showPage()  # página totalmente vacía
    c.save()
    empty_pdf = buf.getvalue()

    with mock.patch("app.documents.extractor._ocr_pdf", return_value=[]):
        with pytest.raises(ValueError, match=r"texto|extraíble"):
            extract_text(empty_pdf)


def test_pdf_with_no_text_falls_back_to_ocr():
    """Si el PDF no tiene texto seleccionable, se intenta OCR antes de fallar."""
    pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas  # type: ignore[import-not-found]

    buf = BytesIO()
    c = canvas.Canvas(buf)
    c.showPage()
    c.save()
    empty_pdf = buf.getvalue()

    with mock.patch(
        "app.documents.extractor._ocr_pdf", return_value=["texto detectado por ocr"]
    ) as mocked_ocr:
        text = extract_text(empty_pdf)

    mocked_ocr.assert_called_once()
    assert "ocr" in text.lower()


# ============================================================
# PPTX
# ============================================================

def _make_pptx_bytes(slide_texts: list[str]) -> bytes:
    pytest.importorskip("pptx", reason="python-pptx no instalado")
    from pptx import Presentation  # type: ignore[import-not-found]

    prs = Presentation()
    layout = prs.slide_layouts[1]  # layout con title + content placeholder
    for slide_text in slide_texts:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_text

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_extract_pptx_reads_slide_text():
    pptx_bytes = _make_pptx_bytes(["Introducción a las derivadas", "Regla de la cadena"])
    text = extract_text(pptx_bytes, mime_type=PPTX_MIME)
    assert "derivadas" in text.lower()
    assert "cadena" in text.lower()


def test_pptx_without_text_or_images_raises():
    pytest.importorskip("pptx", reason="python-pptx no instalado")
    from pptx import Presentation  # type: ignore[import-not-found]

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # layout en blanco, sin placeholders
    buf = BytesIO()
    prs.save(buf)

    with pytest.raises(ValueError, match=r"extraíble"):
        extract_text(buf.getvalue(), mime_type=PPTX_MIME)


# ============================================================
# XLSX
# ============================================================

def _make_xlsx_bytes(rows: list[list]) -> bytes:
    openpyxl = pytest.importorskip("openpyxl", reason="openpyxl no instalado")
    wb = openpyxl.Workbook()
    ws = wb.active
    for row in rows:
        ws.append(row)
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def test_extract_xlsx_reads_cell_values():
    xlsx_bytes = _make_xlsx_bytes([["Tema", "Nota"], ["Límites", 9], ["Derivadas", 8]])
    text = extract_text(xlsx_bytes, mime_type=XLSX_MIME)
    assert "límites" in text.lower()
    assert "derivadas" in text.lower()


def test_empty_xlsx_raises():
    xlsx_bytes = _make_xlsx_bytes([])
    with pytest.raises(ValueError, match=r"extraíbles"):
        extract_text(xlsx_bytes, mime_type=XLSX_MIME)


# ============================================================
# CSV
# ============================================================

def test_extract_csv_joins_rows():
    csv_bytes = "tema,nota\nlímites,9\nderivadas,8\n".encode("utf-8")
    text = extract_text(csv_bytes, mime_type=CSV_MIME)
    assert "límites" in text.lower()
    assert "derivadas" in text.lower()


def test_empty_csv_raises():
    with pytest.raises(ValueError, match=r"extraíbles"):
        extract_text(b"\n\n", mime_type=CSV_MIME)


# ============================================================
# Markdown
# ============================================================

def test_extract_markdown_returns_plain_text():
    md_bytes = "# Título\n\nContenido sobre **derivadas**.".encode("utf-8")
    text = extract_text(md_bytes, mime_type=MD_MIME)
    assert "derivadas" in text.lower()


# ============================================================
# HTML
# ============================================================

def test_extract_html_strips_tags():
    html_bytes = (
        "<html><head><style>body{color:red}</style></head>"
        "<body><h1>Derivadas</h1><script>alert(1)</script>"
        "<p>Contenido de prueba.</p></body></html>"
    ).encode("utf-8")
    text = extract_text(html_bytes, mime_type=HTML_MIME)
    assert "derivadas" in text.lower()
    assert "contenido de prueba" in text.lower()
    assert "alert" not in text.lower()
    assert "color:red" not in text.lower()


def test_html_without_visible_text_raises():
    html_bytes = b"<html><body><script>alert(1)</script></body></html>"
    with pytest.raises(ValueError, match=r"extraíble"):
        extract_text(html_bytes, mime_type=HTML_MIME)


# ============================================================
# Tipo no soportado
# ============================================================

def test_unsupported_mime_type_raises():
    with pytest.raises(ValueError, match="no soportado"):
        extract_text(b"contenido", mime_type="application/zip")
